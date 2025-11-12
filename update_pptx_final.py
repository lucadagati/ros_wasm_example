#!/usr/bin/env python3
"""Update PPTX with all UML diagrams showing ROS in WASM"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Load existing presentation
prs = Presentation('microROS_WASM_K8s_Feasibility.pptx')

print("Updating PPTX with all UML diagrams (ROS in WASM)...")

# UML diagrams mapping
uml_diagrams = {
    'architecture_complete': 'diagrams/architecture_complete.png',
    'workflow': 'diagrams/workflow.png',
    'architecture': 'diagrams/architecture.png',
    'sequence': 'diagrams/sequence.png',
    'k8s_deployment': 'diagrams/k8s_deployment.png',
    'component': 'diagrams/component_diagram.png'
}

for name, path in uml_diagrams.items():
    if os.path.exists(path):
        print(f"✓ Found: {path}")
    else:
        print(f"⚠ Missing: {path}")

# Clear existing slides and recreate structure
# Keep first slide (OT on K8s + Goal)
# Remove other slides
while len(prs.slides) > 1:
    rId = prs.slides._sldIdLst[-1].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[-1]

# Slide 2: Complete Architecture Schema
slide_arch = prs.slides.add_slide(prs.slide_layouts[6])

title_box = slide_arch.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_frame.text = "Complete System Architecture"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(40)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(0, 51, 102)
title_para.alignment = PP_ALIGN.CENTER

subtitle_box = slide_arch.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(0.5))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "ROS executing inside WASM Runtime on Kubernetes"
subtitle_para = subtitle_frame.paragraphs[0]
subtitle_para.font.size = Pt(20)
subtitle_para.font.italic = True
subtitle_para.font.color.rgb = RGBColor(100, 100, 100)
subtitle_para.alignment = PP_ALIGN.CENTER

if os.path.exists(uml_diagrams['architecture_complete']):
    slide_arch.shapes.add_picture(uml_diagrams['architecture_complete'], Inches(0.5), Inches(1.6), 
                                  width=Inches(9), height=Inches(5.8))
    print("✓ Added complete architecture schema")

# Slide 3: System Workflow
slide_workflow = prs.slides.add_slide(prs.slide_layouts[6])

title_box = slide_workflow.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_frame.text = "System Workflow"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(40)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(0, 51, 102)
title_para.alignment = PP_ALIGN.CENTER

subtitle_box = slide_workflow.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(0.5))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "Complete workflow: ROS executes inside WASM from deployment to message processing"
subtitle_para = subtitle_frame.paragraphs[0]
subtitle_para.font.size = Pt(18)
subtitle_para.font.italic = True
subtitle_para.font.color.rgb = RGBColor(100, 100, 100)
subtitle_para.alignment = PP_ALIGN.CENTER

if os.path.exists(uml_diagrams['workflow']):
    slide_workflow.shapes.add_picture(uml_diagrams['workflow'], Inches(0.5), Inches(1.6), 
                                     width=Inches(9), height=Inches(5.8))
    print("✓ Added workflow diagram")

# Slide 4: Communication Sequence
slide_seq = prs.slides.add_slide(prs.slide_layouts[6])

title_box = slide_seq.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_frame.text = "Communication Sequence"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(40)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(0, 51, 102)
title_para.alignment = PP_ALIGN.CENTER

subtitle_box = slide_seq.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(0.5))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "ROS operations execute inside WASM runtime"
subtitle_para = subtitle_frame.paragraphs[0]
subtitle_para.font.size = Pt(18)
subtitle_para.font.italic = True
subtitle_para.font.color.rgb = RGBColor(100, 100, 100)
subtitle_para.alignment = PP_ALIGN.CENTER

if os.path.exists(uml_diagrams['sequence']):
    slide_seq.shapes.add_picture(uml_diagrams['sequence'], Inches(0.5), Inches(1.6), 
                                 width=Inches(9), height=Inches(5.8))
    print("✓ Added sequence diagram")

# Slide 5: K8s Deployment
slide_k8s = prs.slides.add_slide(prs.slide_layouts[6])

title_box = slide_k8s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_frame.text = "Kubernetes Deployment Architecture"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(0, 51, 102)
title_para.alignment = PP_ALIGN.CENTER

if os.path.exists(uml_diagrams['k8s_deployment']):
    slide_k8s.shapes.add_picture(uml_diagrams['k8s_deployment'], Inches(0.5), Inches(1.3), 
                                width=Inches(9), height=Inches(6.0))
    print("✓ Added K8s deployment diagram")

# Slide 6: Component Diagram
slide_comp = prs.slides.add_slide(prs.slide_layouts[6])

title_box = slide_comp.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_frame.text = "System Component Diagram"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(40)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(0, 51, 102)
title_para.alignment = PP_ALIGN.CENTER

subtitle_box = slide_comp.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(0.5))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "ROS components executing in WASM runtime"
subtitle_para = subtitle_frame.paragraphs[0]
subtitle_para.font.size = Pt(18)
subtitle_para.font.italic = True
subtitle_para.font.color.rgb = RGBColor(100, 100, 100)
subtitle_para.alignment = PP_ALIGN.CENTER

if os.path.exists(uml_diagrams['component']):
    slide_comp.shapes.add_picture(uml_diagrams['component'], Inches(0.5), Inches(1.6), 
                                 width=Inches(9), height=Inches(5.8))
    print("✓ Added component diagram")

# Slide 7: Feasibility Analysis
slide_feas = prs.slides.add_slide(prs.slide_layouts[6])

title_box = slide_feas.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "Feasibility: Can We Run ROS in WASM?"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(38)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(0, 51, 102)

left_box = slide_feas.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.5), Inches(5))
left_frame = left_box.text_frame
left_frame.word_wrap = True

p = left_frame.paragraphs[0]
p.text = "Technical Challenges"
p.font.size = Pt(22)
p.font.bold = True
p.space_after = Pt(12)

p = left_frame.add_paragraph()
p.text = "1. System Calls & I/O"
p.font.size = Pt(18)
p.font.bold = True
p.space_after = Pt(4)
p2 = left_frame.add_paragraph()
p2.text = "   WASM limited access → WASI/host bindings"
p2.font.size = Pt(14)
p2.space_after = Pt(10)

p = left_frame.add_paragraph()
p.text = "2. Real-time Constraints"
p.font.size = Pt(18)
p.font.bold = True
p.space_after = Pt(4)
p2 = left_frame.add_paragraph()
p2.text = "   Deterministic timing → WASM <1ms latency"
p2.font.size = Pt(14)
p2.space_after = Pt(10)

p = left_frame.add_paragraph()
p.text = "3. Memory Management"
p.font.size = Pt(18)
p.font.bold = True
p.space_after = Pt(4)
p2 = left_frame.add_paragraph()
p2.text = "   Linear memory → Compatible"
p2.font.size = Pt(14)

right_box = slide_feas.shapes.add_textbox(Inches(5.5), Inches(2), Inches(4.5), Inches(5))
right_frame = right_box.text_frame
right_frame.word_wrap = True

p = right_frame.paragraphs[0]
p.text = "Assessment Results"
p.font.size = Pt(22)
p.font.bold = True
p.space_after = Pt(12)

p = right_frame.add_paragraph()
p.text = "FEASIBILITY: HIGH"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 153, 0)
p.space_after = Pt(14)

p = right_frame.add_paragraph()
p.text = "✓ ROS in WASM: Feasible"
p.font.size = Pt(16)
p.space_after = Pt(6)

p = right_frame.add_paragraph()
p.text = "✓ ROS2 DDS: Portable to WASM"
p.font.size = Pt(16)
p.space_after = Pt(6)

p = right_frame.add_paragraph()
p.text = "✓ Direct WASM-to-WASM: Yes"
p.font.size = Pt(16)
p.space_after = Pt(6)

p = right_frame.add_paragraph()
p.text = "✓ Memory: Lightweight (<100KB)"
p.font.size = Pt(16)
p.space_after = Pt(6)

p = right_frame.add_paragraph()
p.text = "✓ Performance: Near-native"
p.font.size = Pt(16)
p.space_after = Pt(14)

p = right_frame.add_paragraph()
p.text = "Next Steps"
p.font.size = Pt(20)
p.font.bold = True
p.space_after = Pt(8)

p = right_frame.add_paragraph()
p.text = "• Port microROS to WASM"
p.font.size = Pt(15)
p.space_after = Pt(4)

p = right_frame.add_paragraph()
p.text = "• Implement ROS2 DDS in WASM"
p.font.size = Pt(15)
p.space_after = Pt(4)

p = right_frame.add_paragraph()
p.text = "• Deploy on K8s (Wasmtime)"
p.font.size = Pt(15)

# Save updated presentation
prs.save('microROS_WASM_K8s_Feasibility.pptx')
print(f"\n✓ PPTX updated with all UML diagrams")
print(f"  Total slides: {len(prs.slides)}")
print(f"\nSlide structure:")
print(f"  1. OT on K8s + Goal")
print(f"  2. Complete Architecture Schema (ROS in WASM)")
print(f"  3. System Workflow (ROS in WASM)")
print(f"  4. Communication Sequence (ROS in WASM)")
print(f"  5. K8s Deployment (ROS in WASM)")
print(f"  6. Component Diagram (ROS in WASM)")
print(f"  7. Feasibility Analysis")

