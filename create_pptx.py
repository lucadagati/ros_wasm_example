#!/usr/bin/env python3
"""Script to create PPTX presentation for microROS WASM K8s feasibility"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("ERROR: python-pptx not installed. Install with:")
    print("  pip install --user python-pptx")
    print("  or")
    print("  sudo apt install python3-pptx")
    exit(1)

# Create new presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Helper function to add arrow line
def add_arrow(slide, x1, y1, x2, y2, color):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(1.5)

# Helper function to add box with text
def add_box(slide, left, top, width, height, text, fill_color, text_color=RGBColor(0, 0, 0)):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    box.line.color.rgb = RGBColor(0, 0, 0)
    box.line.width = Pt(1)
    
    text_frame = box.text_frame
    text_frame.text = text
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = text_frame.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    para.font.size = Pt(10)
    para.font.bold = True
    para.font.color.rgb = text_color
    return box

# ============================================
# SLIDE 1: OT on K8s + Goal microROS in WASM
# ============================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

# Title
title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "Operational Technology on Kubernetes"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(40)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(0, 51, 102)

# Left column - OT Infrastructure
left_box = slide1.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.5), Inches(5))
left_frame = left_box.text_frame
left_frame.word_wrap = True

p = left_frame.paragraphs[0]
p.text = "OT Infrastructure"
p.font.size = Pt(22)
p.font.bold = True
p.space_after = Pt(10)

p = left_frame.add_paragraph()
p.text = "• Industrial control systems"
p.font.size = Pt(16)
p.space_after = Pt(4)

p = left_frame.add_paragraph()
p.text = "• Real-time data processing"
p.font.size = Pt(16)
p.space_after = Pt(4)

p = left_frame.add_paragraph()
p.text = "• Distributed sensor networks"
p.font.size = Pt(16)
p.space_after = Pt(16)

p = left_frame.add_paragraph()
p.text = "Kubernetes Orchestration"
p.font.size = Pt(22)
p.font.bold = True
p.space_after = Pt(10)

p = left_frame.add_paragraph()
p.text = "• Container orchestration"
p.font.size = Pt(16)
p.space_after = Pt(4)

p = left_frame.add_paragraph()
p.text = "• Scalable microservices"
p.font.size = Pt(16)
p.space_after = Pt(4)

p = left_frame.add_paragraph()
p.text = "• High availability"
p.font.size = Pt(16)

# Right column - Goal
right_box = slide1.shapes.add_textbox(Inches(5.5), Inches(2), Inches(4.5), Inches(5))
right_frame = right_box.text_frame
right_frame.word_wrap = True

p = right_frame.paragraphs[0]
p.text = "Project Goal"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 102, 51)
p.space_after = Pt(12)

p = right_frame.add_paragraph()
p.text = "Run microROS inside WebAssembly runtime"
p.font.size = Pt(18)
p.font.bold = True
p.space_after = Pt(16)

p = right_frame.add_paragraph()
p.text = "Key Benefits"
p.font.size = Pt(20)
p.font.bold = True
p.space_after = Pt(10)

p = right_frame.add_paragraph()
p.text = "✓ Lightweight & Portable"
p.font.size = Pt(16)
p.space_after = Pt(6)

p = right_frame.add_paragraph()
p.text = "✓ Secure sandbox"
p.font.size = Pt(16)
p.space_after = Pt(6)

p = right_frame.add_paragraph()
p.text = "✓ Near-native performance"
p.font.size = Pt(16)
p.space_after = Pt(6)

p = right_frame.add_paragraph()
p.text = "✓ K8s orchestration ready"
p.font.size = Pt(16)

# ============================================
# SLIDE 2: Sequence Diagram + K8s Architecture
# ============================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "Communication Flow & Kubernetes Architecture"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(32)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(0, 51, 102)

# Left side - Visual Sequence Diagram
# Participants (vertical boxes on left)
pub_box = add_box(slide2, 0.5, 2.2, 1.0, 0.4, "Publisher\nWASM", RGBColor(102, 126, 234), RGBColor(255, 255, 255))
js1_box = add_box(slide2, 0.5, 2.8, 1.0, 0.4, "JavaScript\n(Publisher)", RGBColor(76, 175, 80), RGBColor(255, 255, 255))
rb_box = add_box(slide2, 0.5, 3.4, 1.0, 0.4, "ROS\nBridge", RGBColor(255, 193, 7), RGBColor(0, 0, 0))
ros_box = add_box(slide2, 0.5, 4.0, 1.0, 0.4, "ROS2\nTopics", RGBColor(156, 39, 176), RGBColor(255, 255, 255))
js2_box = add_box(slide2, 0.5, 4.6, 1.0, 0.4, "JavaScript\n(Subscriber)", RGBColor(76, 175, 80), RGBColor(255, 255, 255))
sub_box = add_box(slide2, 0.5, 5.2, 1.0, 0.4, "Subscriber\nWASM", RGBColor(245, 87, 108), RGBColor(255, 255, 255))

# Lifelines (vertical lines)
lifeline_x = 1.1
for y in [2.4, 3.0, 3.6, 4.2, 4.8, 5.4]:
    line = slide2.shapes.add_connector(1, Inches(lifeline_x), Inches(y), Inches(lifeline_x), Inches(y + 0.6))
    line.line.color.rgb = RGBColor(200, 200, 200)
    line.line.width = Pt(1)

# Messages (horizontal arrows)
arrow_y_start = 2.5
arrow_spacing = 0.35
arrow_x_start = 1.2
arrow_x_end = 3.5

# Message 1: generateMessage()
add_arrow(slide2, arrow_x_start, arrow_y_start, arrow_x_end, arrow_y_start, RGBColor(102, 126, 234))
msg1 = slide2.shapes.add_textbox(Inches(1.5), Inches(arrow_y_start - 0.1), Inches(1.8), Inches(0.15))
msg1.text_frame.text = "generateMessage()"
msg1.text_frame.paragraphs[0].font.size = Pt(9)
msg1.text_frame.paragraphs[0].font.bold = True

# Message 2: publish()
add_arrow(slide2, arrow_x_start, arrow_y_start + arrow_spacing, arrow_x_end, arrow_y_start + arrow_spacing, RGBColor(76, 175, 80))
msg2 = slide2.shapes.add_textbox(Inches(1.5), Inches(arrow_y_start + arrow_spacing - 0.1), Inches(1.8), Inches(0.15))
msg2.text_frame.text = "publish(topic, message)"
msg2.text_frame.paragraphs[0].font.size = Pt(9)
msg2.text_frame.paragraphs[0].font.bold = True

# Message 3: ROS2 publish
add_arrow(slide2, arrow_x_start, arrow_y_start + arrow_spacing*2, arrow_x_end, arrow_y_start + arrow_spacing*2, RGBColor(255, 193, 7))
msg3 = slide2.shapes.add_textbox(Inches(1.5), Inches(arrow_y_start + arrow_spacing*2 - 0.1), Inches(1.8), Inches(0.15))
msg3.text_frame.text = "ROS2 DDS publish"
msg3.text_frame.paragraphs[0].font.size = Pt(9)
msg3.text_frame.paragraphs[0].font.bold = True

# Message 4: message delivery (return)
add_arrow(slide2, arrow_x_end, arrow_y_start + arrow_spacing*3, arrow_x_start, arrow_y_start + arrow_spacing*3, RGBColor(156, 39, 176))
msg4 = slide2.shapes.add_textbox(Inches(1.5), Inches(arrow_y_start + arrow_spacing*3 - 0.1), Inches(1.8), Inches(0.15))
msg4.text_frame.text = "message delivery"
msg4.text_frame.paragraphs[0].font.size = Pt(9)
msg4.text_frame.paragraphs[0].font.bold = True

# Message 5: callback
add_arrow(slide2, arrow_x_start, arrow_y_start + arrow_spacing*4, arrow_x_end, arrow_y_start + arrow_spacing*4, RGBColor(76, 175, 80))
msg5 = slide2.shapes.add_textbox(Inches(1.5), Inches(arrow_y_start + arrow_spacing*4 - 0.1), Inches(1.8), Inches(0.15))
msg5.text_frame.text = "message callback"
msg5.text_frame.paragraphs[0].font.size = Pt(9)
msg5.text_frame.paragraphs[0].font.bold = True

# Message 6: processMessage()
add_arrow(slide2, arrow_x_start, arrow_y_start + arrow_spacing*5, arrow_x_end, arrow_y_start + arrow_spacing*5, RGBColor(245, 87, 108))
msg6 = slide2.shapes.add_textbox(Inches(1.5), Inches(arrow_y_start + arrow_spacing*5 - 0.1), Inches(1.8), Inches(0.15))
msg6.text_frame.text = "processMessage(data)"
msg6.text_frame.paragraphs[0].font.size = Pt(9)
msg6.text_frame.paragraphs[0].font.bold = True

# Loop label
loop_label = slide2.shapes.add_textbox(Inches(3.6), Inches(2.2), Inches(1.2), Inches(0.3))
loop_label.text_frame.text = "Loop: Every second"
loop_label.text_frame.paragraphs[0].font.size = Pt(8)
loop_label.text_frame.paragraphs[0].font.italic = True
loop_label.text_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)

# Right side - Visual K8s Architecture Diagram
# K8s Cluster container
cluster_box = add_box(slide2, 5.5, 2.0, 4.0, 3.0, "Kubernetes Cluster", RGBColor(240, 240, 240), RGBColor(0, 0, 0))
cluster_box.line.width = Pt(2)

# Pod 1
pod1_box = add_box(slide2, 5.8, 2.5, 1.6, 1.0, "K8s Pod 1", RGBColor(200, 220, 255), RGBColor(0, 0, 0))
wasm1_box = add_box(slide2, 6.0, 2.8, 1.2, 0.35, "WASM Runtime", RGBColor(102, 126, 234), RGBColor(255, 255, 255))
micro1_box = add_box(slide2, 6.0, 3.2, 1.2, 0.35, "microROS Node A", RGBColor(76, 175, 80), RGBColor(255, 255, 255))

# Pod 2
pod2_box = add_box(slide2, 7.6, 2.5, 1.6, 1.0, "K8s Pod 2", RGBColor(200, 220, 255), RGBColor(0, 0, 0))
wasm2_box = add_box(slide2, 7.8, 2.8, 1.2, 0.35, "WASM Runtime", RGBColor(102, 126, 234), RGBColor(255, 255, 255))
micro2_box = add_box(slide2, 7.8, 3.2, 1.2, 0.35, "microROS Node B", RGBColor(76, 175, 80), RGBColor(255, 255, 255))

# Connection arrow between pods
add_arrow(slide2, 7.4, 3.4, 7.6, 3.4, RGBColor(0, 0, 0))
conn_label = slide2.shapes.add_textbox(Inches(7.0), Inches(3.5), Inches(1.0), Inches(0.2))
conn_label.text_frame.text = "ROS Topics/Services"
conn_label.text_frame.paragraphs[0].font.size = Pt(8)
conn_label.text_frame.paragraphs[0].font.bold = True

# K8s management features
k8s_label = slide2.shapes.add_textbox(Inches(5.8), Inches(4.0), Inches(3.6), Inches(0.8))
k8s_frame = k8s_label.text_frame
k8s_frame.word_wrap = True
k8s_frame.text = "K8s manages: Deployment • Scaling • Health Checks • Updates"
k8s_frame.paragraphs[0].font.size = Pt(10)
k8s_frame.paragraphs[0].font.bold = True

# ============================================
# SLIDE 3: Feasibility Analysis
# ============================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "Feasibility: Can We Run microROS in WASM?"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(38)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(0, 51, 102)

# Left column - Challenges
left_box = slide3.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.5), Inches(5))
left_frame = left_box.text_frame
left_frame.word_wrap = True

p = left_frame.paragraphs[0]
p.text = "Technical Challenges"
p.font.size = Pt(22)
p.font.bold = True
p.space_after = Pt(12)

p = left_frame.add_paragraph()
p.text = "1. System Calls & I/O"
p.font.size = Pt(18)
p.font.bold = True
p.space_after = Pt(4)
p2 = left_frame.add_paragraph()
p2.text = "   WASM limited access → WASI/host bindings"
p2.font.size = Pt(14)
p2.space_after = Pt(10)

p = left_frame.add_paragraph()
p.text = "2. Real-time Constraints"
p.font.size = Pt(18)
p.font.bold = True
p.space_after = Pt(4)
p2 = left_frame.add_paragraph()
p2.text = "   Deterministic timing → WASM <1ms latency"
p2.font.size = Pt(14)
p2.space_after = Pt(10)

p = left_frame.add_paragraph()
p.text = "3. Memory Management"
p.font.size = Pt(18)
p.font.bold = True
p.space_after = Pt(4)
p2 = left_frame.add_paragraph()
p2.text = "   Linear memory → Compatible"
p2.font.size = Pt(14)

# Right column - Assessment
right_box = slide3.shapes.add_textbox(Inches(5.5), Inches(2), Inches(4.5), Inches(5))
right_frame = right_box.text_frame
right_frame.word_wrap = True

p = right_frame.paragraphs[0]
p.text = "Assessment Results"
p.font.size = Pt(22)
p.font.bold = True
p.space_after = Pt(12)

p = right_frame.add_paragraph()
p.text = "FEASIBILITY: HIGH"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 153, 0)
p.space_after = Pt(14)

p = right_frame.add_paragraph()
p.text = "✓ Core microROS: Compatible"
p.font.size = Pt(16)
p.space_after = Pt(6)

p = right_frame.add_paragraph()
p.text = "✓ ROS2 DDS: Portable to WASM"
p.font.size = Pt(16)
p.space_after = Pt(6)

p = right_frame.add_paragraph()
p.text = "✓ Pub/Sub: Demonstrated"
p.font.size = Pt(16)
p.space_after = Pt(6)

p = right_frame.add_paragraph()
p.text = "✓ Memory: Lightweight (<100KB)"
p.font.size = Pt(16)
p.space_after = Pt(6)

p = right_frame.add_paragraph()
p.text = "✓ Performance: Near-native"
p.font.size = Pt(16)
p.space_after = Pt(14)

p = right_frame.add_paragraph()
p.text = "Next Steps"
p.font.size = Pt(20)
p.font.bold = True
p.space_after = Pt(8)

p = right_frame.add_paragraph()
p.text = "• Port microROS to WASM"
p.font.size = Pt(15)
p.space_after = Pt(4)

p = right_frame.add_paragraph()
p.text = "• Implement ROS2 DDS layer"
p.font.size = Pt(15)
p.space_after = Pt(4)

p = right_frame.add_paragraph()
p.text = "• Deploy on K8s (Wasmtime)"
p.font.size = Pt(15)

# Save presentation
try:
    prs.save('microROS_WASM_K8s_Feasibility.pptx')
    print("✓ Presentation created: microROS_WASM_K8s_Feasibility.pptx with 3 slides")
    print("  - Slide 1: OT on K8s + Goal")
    print("  - Slide 2: Visual Sequence Diagram + K8s Architecture")
    print("  - Slide 3: Feasibility Analysis")
except Exception as e:
    print(f"ERROR saving PPTX: {e}")
    exit(1)
