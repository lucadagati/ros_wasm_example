#!/usr/bin/env python3
"""Update PPTX with architecture schema and workflow"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Load existing presentation
prs = Presentation('microROS_WASM_K8s_Feasibility.pptx')

print("Updating PPTX with architecture schema and workflow...")

# Check if diagrams exist
diagrams = {
    'architecture_complete': 'diagrams/architecture_complete.png',
    'workflow': 'diagrams/workflow.png',
}

for name, path in diagrams.items():
    if os.path.exists(path):
        print(f"✓ Found: {path}")
    else:
        print(f"⚠ Missing: {path}")

# Create new slide for Complete Architecture Schema
slide_arch = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_box = slide_arch.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_frame.text = "Complete System Architecture"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(40)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(0, 51, 102)
title_para.alignment = PP_ALIGN.CENTER

# Subtitle
subtitle_box = slide_arch.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(0.5))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "microROS in WASM Runtime on Kubernetes"
subtitle_para = subtitle_frame.paragraphs[0]
subtitle_para.font.size = Pt(20)
subtitle_para.font.italic = True
subtitle_para.font.color.rgb = RGBColor(100, 100, 100)
subtitle_para.alignment = PP_ALIGN.CENTER

# Add complete architecture diagram
if os.path.exists(diagrams['architecture_complete']):
    slide_arch.shapes.add_picture(diagrams['architecture_complete'], Inches(0.5), Inches(1.6), 
                                  width=Inches(9), height=Inches(5.8))
    print("✓ Added complete architecture schema")

# Create new slide for Workflow
slide_workflow = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_box = slide_workflow.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_frame.text = "System Workflow"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(40)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(0, 51, 102)
title_para.alignment = PP_ALIGN.CENTER

# Subtitle
subtitle_box = slide_workflow.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(0.5))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "Complete workflow from deployment to message processing"
subtitle_para = subtitle_frame.paragraphs[0]
subtitle_para.font.size = Pt(18)
subtitle_para.font.italic = True
subtitle_para.font.color.rgb = RGBColor(100, 100, 100)
subtitle_para.alignment = PP_ALIGN.CENTER

# Add workflow diagram
if os.path.exists(diagrams['workflow']):
    slide_workflow.shapes.add_picture(diagrams['workflow'], Inches(0.5), Inches(1.6), 
                                     width=Inches(9), height=Inches(5.8))
    print("✓ Added workflow diagram")

# Save updated presentation
prs.save('microROS_WASM_K8s_Feasibility.pptx')
print(f"\n✓ PPTX updated with architecture schema and workflow")
print(f"  Total slides: {len(prs.slides)}")
print(f"\nSlide structure:")
print(f"  1. OT on K8s + Goal")
print(f"  2. Communication Flow & Sequence")
print(f"  3. K8s Deployment")
print(f"  4. Component & Class Diagrams")
print(f"  5. Feasibility Analysis")
print(f"  6. Complete Architecture Schema (NEW)")
print(f"  7. System Workflow (NEW)")
