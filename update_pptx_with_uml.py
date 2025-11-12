#!/usr/bin/env python3
"""Update PPTX with formal UML diagrams from PlantUML"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Load existing presentation
prs = Presentation('microROS_WASM_K8s_Feasibility.pptx')

# UML diagrams mapping
uml_diagrams = {
    'architecture': 'diagrams/architecture.png',
    'sequence': 'diagrams/sequence.png',
    'k8s_deployment': 'diagrams/k8s_deployment.png',
    'component': 'diagrams/component_diagram.png',
    'class': 'diagrams/class_diagram.png'
}

print("Updating PPTX with UML diagrams...")

# Update Slide 2 - Replace with UML diagrams
slide2 = prs.slides[1]  # Second slide

# Clear existing content (keep title)
shapes_to_remove = []
for shape in slide2.shapes:
    if shape.shape_type == 17:  # Text box
        if "Communication Flow" not in shape.text_frame.text[:30]:
            shapes_to_remove.append(shape)
    elif shape.shape_type != 6:  # Not a picture
        if shape.shape_type in [1, 25]:  # AutoShape or Connector
            shapes_to_remove.append(shape)

for shape in shapes_to_remove:
    try:
        sp = shape._element
        sp.getparent().remove(sp)
    except:
        pass

# Add UML Architecture diagram (left side)
if os.path.exists(uml_diagrams['architecture']):
    slide2.shapes.add_picture(uml_diagrams['architecture'], Inches(0.5), Inches(1.8), 
                             width=Inches(4.5), height=Inches(3.2))
    print("✓ Added UML architecture diagram")

# Add UML Sequence diagram (right side)
if os.path.exists(uml_diagrams['sequence']):
    slide2.shapes.add_picture(uml_diagrams['sequence'], Inches(5.5), Inches(1.8), 
                             width=Inches(4.5), height=Inches(3.2))
    print("✓ Added UML sequence diagram")

# Update or create K8s deployment slide
if len(prs.slides) > 3:
    slide_k8s = prs.slides[3]  # Existing K8s slide
    # Clear it
    for shape in list(slide_k8s.shapes):
        if shape.shape_type != 17 or "Kubernetes" not in shape.text_frame.text[:20]:
            try:
                sp = shape._element
                sp.getparent().remove(sp)
            except:
                pass
else:
    slide_k8s = prs.slides.add_slide(prs.slide_layouts[6])

# Title for K8s slide
title_box = slide_k8s.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "Kubernetes Deployment Architecture"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(0, 51, 102)

# Add UML K8s deployment diagram
if os.path.exists(uml_diagrams['k8s_deployment']):
    slide_k8s.shapes.add_picture(uml_diagrams['k8s_deployment'], Inches(0.5), Inches(1.8), 
                                width=Inches(9), height=Inches(5.5))
    print("✓ Added UML K8s deployment diagram")

# Create new slide for Component and Class diagrams
slide_detailed = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_box = slide_detailed.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "System Component & Class Diagrams"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(0, 51, 102)

# Add Component diagram (left)
if os.path.exists(uml_diagrams['component']):
    slide_detailed.shapes.add_picture(uml_diagrams['component'], Inches(0.5), Inches(1.8), 
                                     width=Inches(4.5), height=Inches(3.2))
    print("✓ Added UML component diagram")

# Add Class diagram (right)
if os.path.exists(uml_diagrams['class']):
    slide_detailed.shapes.add_picture(uml_diagrams['class'], Inches(5.5), Inches(1.8), 
                                     width=Inches(4.5), height=Inches(3.2))
    print("✓ Added UML class diagram")

# Save updated presentation
prs.save('microROS_WASM_K8s_Feasibility.pptx')
print("\n✓ PPTX updated with formal UML diagrams")
print(f"  Total slides: {len(prs.slides)}")

