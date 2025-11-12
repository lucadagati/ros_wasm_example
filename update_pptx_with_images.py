#!/usr/bin/env python3
"""Update PPTX with UML diagram images"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Load existing presentation
prs = Presentation('microROS_WASM_K8s_Feasibility.pptx')

# Check if diagrams exist
diagrams = {
    'architecture': 'diagrams/architecture_diagram.png',
    'sequence': 'diagrams/sequence_diagram.png',
    'k8s': 'diagrams/k8s_deployment_diagram.png'
}

for name, path in diagrams.items():
    if not os.path.exists(path):
        print(f"WARNING: {path} not found")
    else:
        print(f"✓ Found: {path}")

# Update Slide 2 - Replace sequence diagram and K8s architecture with images
slide2 = prs.slides[1]  # Second slide (index 1)

# Remove existing shapes (except title)
shapes_to_remove = []
for shape in slide2.shapes:
    if shape.shape_type == 17:  # Text box
        if "Communication Flow" not in shape.text_frame.text[:30]:
            shapes_to_remove.append(shape)
    elif shape.shape_type == 1:  # AutoShape
        shapes_to_remove.append(shape)
    elif shape.shape_type == 25:  # Connector
        shapes_to_remove.append(shape)

for shape in shapes_to_remove:
    sp = shape._element
    sp.getparent().remove(sp)

# Add architecture diagram (left side)
if os.path.exists(diagrams['architecture']):
    slide2.shapes.add_picture(diagrams['architecture'], Inches(0.5), Inches(1.8), 
                             width=Inches(4.5), height=Inches(3.0))
    print("✓ Added architecture diagram to slide 2")

# Add sequence diagram (right side, top)
if os.path.exists(diagrams['sequence']):
    slide2.shapes.add_picture(diagrams['sequence'], Inches(5.5), Inches(1.8), 
                             width=Inches(4.5), height=Inches(3.0))
    print("✓ Added sequence diagram to slide 2")

# Add K8s deployment diagram to a new slide or replace content
# Create new slide for K8s deployment
slide_k8s = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_box = slide_k8s.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "Kubernetes Deployment Architecture"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(0, 51, 102)

# Add K8s diagram
if os.path.exists(diagrams['k8s']):
    slide_k8s.shapes.add_picture(diagrams['k8s'], Inches(0.5), Inches(1.8), 
                                width=Inches(9), height=Inches(5.5))
    print("✓ Added K8s deployment diagram to new slide")

# Save updated presentation
prs.save('microROS_WASM_K8s_Feasibility.pptx')
print("\n✓ PPTX updated with UML diagrams")

